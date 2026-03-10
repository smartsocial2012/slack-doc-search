import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import pandas as pd  # Excel, CSV 처리를 위해 추가
import os
import zipfile
import xml.etree.ElementTree as ET

def extract_text_from_pdf(file_path):
    text = ""
    try:
        with fitz.open(file_path) as doc:
            # 1. 문서가 암호화되어 있는지 체크
            if doc.is_encrypted:
                return "실패: 암호화된 PDF 파일입니다."

            for page in doc:
                # 2. 'blocks' 모드를 사용하면 텍스트 묶음 단위로 더 정확하게 읽어옵니다.
                blocks = page.get_text("blocks")
                for b in blocks:
                    # b[4]가 실제 텍스트 내용입니다.
                    text += b[4] + "\n"
            
            # 3. 추출된 텍스트가 너무 적다면? (이미지 PDF일 가능성)
            if len(text.strip()) < 10:
                return "실패: 드래그는 가능하나 텍스트 레이어가 비어있거나 이미지 형태입니다."
                
        return text
    except Exception as e:
        return f"실패: PDF 읽기 오류 ({str(e)})"
    


def extract_text_from_hwpx(file_path):
    """라이브러리 없이 HWPX 내부 XML에서 직접 텍스트를 추출합니다."""
    text = []
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            # HWPX 내부의 실제 본문 데이터는 Contents/section0.xml 등에 들어있습니다.
            content_files = [f for f in zf.namelist() if f.startswith('Contents/section')]
            
            for content_file in content_files:
                with zf.open(content_file) as f:
                    tree = ET.parse(f)
                    root = tree.getroot()
                    # XML 태그 내의 모든 텍스트 요소를 수집
                    for elem in root.iter():
                        if elem.text:
                            text.append(elem.text)
                            
        return "\n".join(text)
    except Exception as e:
        print(f"❌ HWPX 직접 추출 에러: {e}")
        return ""

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

# --- [추가] Excel 추출 로직 ---
def extract_text_from_xlsx(file_path):
    # 모든 시트를 읽어옵니다.
    df_dict = pd.read_excel(file_path, sheet_name=None)
    combined_text = []
    for sheet_name, df in df_dict.items():
        combined_text.append(f"--- 시트명: {sheet_name} ---")
        # 데이터를 CSV 형태의 텍스트로 변환 (구조 유지)
        combined_text.append(df.to_csv(index=False))
    return "\n".join(combined_text)

# --- [추가] CSV 추출 로직 ---
def extract_text_from_csv(file_path):
    # 인코딩 문제가 자주 발생하므로 utf-8-sig(BOM 대응) 사용
    df = pd.read_csv(file_path, encoding='utf-8-sig')
    return df.to_csv(index=False)

def get_text_from_file(file_path):
    ext = file_path.split('.')[-1].lower()
    print(f"🔍 텍스트 추출 중... [{ext.upper()}] : {os.path.basename(file_path)}")
    
    if ext == 'pdf':
        return extract_text_from_pdf(file_path)
    elif ext == 'pptx':
        return extract_text_from_pptx(file_path)
    elif ext == 'docx':
        return extract_text_from_docx(file_path)
    elif ext == 'hwpx':
        return extract_text_from_hwpx(file_path)
    elif ext == 'xlsx':
        return extract_text_from_xlsx(file_path)
    elif ext == 'csv':
        return extract_text_from_csv(file_path)
    elif ext in ['txt', 'md']:
        try:
            # utf-8로 시도하되, 실패하면 한글 인코딩(cp949)까지 고려
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='cp949', errors='ignore') as f:
                return f.read()
    else:
        return None